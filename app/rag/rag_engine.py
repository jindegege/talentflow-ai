# 文件路径：app/core/rag_engine.py
# 作用：RAG 核心逻辑 - 全格式文档加载、清洗、切片（含 OCR 支持）

import os
from typing import List, Dict
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import (
    TextLoader,
)

import docx
from pptx import Presentation

from bs4 import BeautifulSoup

# --- 新增：PDF 处理相关库 ---
import pdfplumber # 用于提取文字版 PDF 和检测页面
from pdf2image import convert_from_path # 用于将 PDF 转为图片
import pytesseract # 调用 Tesseract OCR 引擎

# --- 1. 文档加载器策略 ---

def _load_pdf(file_path: str):
    """
    PDF 解析增强方案：
    1. 优先尝试提取文字（针对文字版 PDF）。
    2. 如果文字极少（针对扫描版 PDF），自动触发 OCR 流程。
    """
    documents = []
    
    # 1. 尝试提取文字
    with pdfplumber.open(file_path) as pdf:
        full_text = ""
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                full_text += text + "\n"
        
        # 2. 判定是否为扫描版
        # 策略：如果提取出的文字非常少（例如少于 50 个字符），认为是扫描版
        if len(full_text.strip()) < 50:
            print(f"检测到 '{file_path}' 可能是扫描版 PDF，正在启动 OCR 识别...")
            documents = _ocr_pdf_pipeline(file_path)
        else:
            # 正常文字版，直接返回
            from langchain_core.documents import Document
            documents.append(Document(page_content=full_text, metadata={"source": file_path}))

    return documents

def _ocr_pdf_pipeline(file_path: str):
    """
    OCR 专用流水线：
    PDF -> 图片 -> OCR 识别 -> 文本
    """
    docs = []
    try:
        # 将 PDF 每一页转为图片 (dpi=300 保证识别清晰度)
        images = convert_from_path(file_path, dpi=300)
        
        full_text = ""
        for i, image in enumerate(images):
            # 使用 pytesseract 识别图片文字
            # lang='chi_sim+eng' 表示同时支持简体中文和英文
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            if text.strip():
                full_text += f"[第 {i+1} 页]\n{text}\n\n"
        
        from langchain_core.documents import Document
        docs.append(Document(page_content=full_text, metadata={"source": file_path}))
        
    except Exception as e:
        print(f"OCR 识别失败: {e}")
        # 如果 OCR 失败，至少返回一个空文档避免程序崩溃
        from langchain_core.documents import Document
        docs.append(Document(page_content="OCR 识别失败，无法提取内容", metadata={"source": file_path}))
        
    return docs

def _load_text(file_path: str):
    """TXT 直接读取"""
    loader = TextLoader(file_path, encoding="utf-8")
    return loader.load()

def _load_docx_pptx(file_path: str):
    """
    Office 文档解析（轻量级方案）：
    不再依赖 Unstructured，直接使用原生库解析，稳定且无警告。
    """
    from langchain_core.documents import Document
    text_content = ""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".docx":
            # 处理 Word 文档
            doc = docx.Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs]
            text_content = "\n".join(paragraphs)
            
        elif ext == ".pptx":
            # 处理 PPT 文档
            prs = Presentation(file_path)
            slide_texts = []
            for i, slide in enumerate(prs.slides):
                slide_text = f"--- 第 {i+1} 页幻灯片 ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                slide_texts.append(slide_text)
            text_content = "\n".join(slide_texts)
            
        # 返回 LangChain 兼容的 Document 对象
        return [Document(page_content=text_content, metadata={"source": file_path})]

    except Exception as e:
        print(f"Office 文档解析失败: {e}")
        return []

def _load_html(file_path: str):
    """网页内容提取：清洗 HTML 标签，去除噪声"""
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
    
    # 去除噪声标签
    for element in soup.find_all(['script', 'style', 'nav', 'header', 'footer', 'aside']):
        element.decompose()
    
    text = soup.get_text(separator="\n", strip=True)
    
    from langchain_core.documents import Document
    return [Document(page_content=text, metadata={"source": file_path})]

# --- 2. 核心流水线 ---
def process_file_pipeline(file_path: str) -> List[Dict]:
    """
    RAG 数据流水线：加载 -> 清洗 -> 切片
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    # 根据文件后缀选择加载策略
    if ext == ".pdf":
        documents = _load_pdf(file_path)
    elif ext == ".txt":
        documents = _load_text(file_path)
    elif ext in [".docx", ".pptx"]:
        documents = _load_docx_pptx(file_path)
    elif ext == ".html":
        documents = _load_html(file_path)
    else:
        # 默认回退到文本加载
        documents = _load_text(file_path)

    # --- 3. 文本切片 (Chunking) ---
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        separators=["\n\n", "\n", "。", "！", "？", "；", " ", ".", "!", "?", ""]
    )
    
    chunks = text_splitter.split_documents(documents)
    
    # --- 4. 格式化输出 ---
    result = []
    for i, chunk in enumerate(chunks):
        result.append({
            "chunk_id": i,
            "content": chunk.page_content.strip(),
            "metadata": {
                "source": chunk.metadata.get("source", file_path),
                "page": chunk.metadata.get("page", 0) + 1,
                "type": ext[1:].upper()
            }
        })
        
    return result